"""
generate_presentation.py — Produces reports/Presentation.pptx
Uses python-pptx to create a 12-slide professional presentation.
"""
import sys
from pathlib import Path

BASE_DIR      = Path(r"c:\Users\Dell\Desktop\Internship_coding")
REPORTS_DIR   = BASE_DIR / "reports"
DATA_DIR      = BASE_DIR / "data" / "processed"
DASHBOARD_DIR = BASE_DIR / "dashboard"

try:
    import pandas as pd
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Cm
    print("python-pptx imports OK")
except ImportError as e:
    print(f"Import error: {e}")
    sys.exit(1)

# ── Color Constants ───────────────────────────────────────────────────────────
C_DARK    = RGBColor(0x0F, 0x17, 0x2A)   # #0F172A
C_SLATE   = RGBColor(0x1E, 0x29, 0x3B)   # #1E293B
C_EMERALD = RGBColor(0x10, 0xB9, 0x81)   # #10B981
C_BLUE    = RGBColor(0x3B, 0x82, 0xF6)   # #3B82F6
C_AMBER   = RGBColor(0xF5, 0x9E, 0x0B)   # #F59E0B
C_RED     = RGBColor(0xEF, 0x44, 0x44)   # #EF4444
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_DIM     = RGBColor(0x94, 0xA3, 0xB8)   # #94A3B8
C_LIGHT   = RGBColor(0xE2, 0xE8, 0xF0)   # #E2E8F0

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h, font_size=18, bold=False,
             color=C_WHITE, align=PP_ALIGN.LEFT, italic=False):
    txbox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txbox


def set_bg(slide, color=C_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_line(slide, l, t, w, color=C_EMERALD, thickness=Pt(2)):
    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR_TYPE.STRAIGHT
        Inches(l), Inches(t), Inches(l + w), Inches(t)
    )
    connector.line.color.rgb = color
    connector.line.width = thickness
    return connector


def add_bullet_box(slide, bullets, l, t, w, h, icon_color=C_EMERALD):
    from pptx.util import Pt
    txbox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run_icon = p.add_run()
        run_icon.text = "▸ "
        run_icon.font.color.rgb = icon_color
        run_icon.font.size = Pt(14)
        run_icon.font.name = "Calibri"
        run = p.add_run()
        run.text = b
        run.font.color.rgb = C_LIGHT
        run.font.size = Pt(13)
        run.font.name = "Calibri"


def add_kpi_box(slide, value, label, l, t, w=1.8, h=1.4,
                val_color=C_EMERALD, bg=C_SLATE):
    add_rect(slide, l, t, w, h, bg)
    add_text(slide, value, l, t + 0.15, w, 0.7,
             font_size=24, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    add_text(slide, label, l, t + 0.85, w, 0.45,
             font_size=10, color=C_DIM, align=PP_ALIGN.CENTER)


def slide_title_block(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.2, C_SLATE)
    add_text(slide, title, 0.5, 0.15, 12, 0.7, font_size=28, bold=True, color=C_EMERALD)
    if subtitle:
        add_text(slide, subtitle, 0.5, 0.85, 12, 0.35, font_size=13, color=C_DIM)
    # Emerald underline
    add_line(slide, 0.5, 1.15, 12.3)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS (12 slides)
# ══════════════════════════════════════════════════════════════════════════════

def slide_01_cover(prs):
    """Slide 1: Title / Cover."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide)
    add_rect(slide, 0, 0, 13.33, 7.5, C_DARK)
    add_rect(slide, 0, 0, 13.33, 3, C_SLATE)
    add_text(slide, "BLUESTOCK FINTECH",
             0.5, 0.4, 12.33, 1.0, font_size=40, bold=True, color=C_EMERALD, align=PP_ALIGN.CENTER)
    add_text(slide, "Mutual Fund Analytics — Capstone Project",
             0.5, 1.4, 12.33, 0.7, font_size=22, color=C_LIGHT, align=PP_ALIGN.CENTER)
    add_text(slide, "Final Presentation",
             0.5, 2.0, 12.33, 0.6, font_size=17, color=C_DIM, align=PP_ALIGN.CENTER, italic=True)
    add_line(slide, 2, 2.9, 9.33)
    # KPI boxes
    kpis = [("40", "Schemes Analysed", 1.2),
            ("51K+", "NAV Data Points", 3.4),
            ("32K+", "Transactions", 5.6),
            ("10", "Raw Datasets", 7.8),
            ("7", "Analysis Days", 10.0)]
    for val, lbl, x in kpis:
        add_kpi_box(slide, val, lbl, x, 3.4, w=2.0, h=1.4)
    add_text(slide, "Prepared by: Kunal Gupta  |  July 2026",
             0.5, 5.5, 12.33, 0.5, font_size=12, color=C_DIM, align=PP_ALIGN.CENTER)
    add_text(slide, "Internship @ Bluestock Fintech  |  AI/ML Data Analytics Track",
             0.5, 5.9, 12.33, 0.5, font_size=11, color=C_DIM, align=PP_ALIGN.CENTER)


def slide_02_problem(prs):
    """Slide 2: Problem Statement & Objectives."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_title_block(slide, "Problem Statement & Objectives", "Why this project matters")
    # Problems column
    add_text(slide, "Key Problems", 0.5, 1.5, 5.5, 0.5,
             font_size=16, bold=True, color=C_EMERALD)
    problems = [
        "Data Fragmentation: NAV, AUM, SIP data scattered across sources",
        "No unified performance comparison across 40+ schemes",
        "Missing benchmark tracking (Nifty 50/100/BSE SmallCap)",
        "Investor behavior blindspot (32K+ transactions unexplored)",
        "Static reporting — no interactive dashboards for stakeholders",
    ]
    add_bullet_box(slide, problems, 0.5, 2.1, 5.8, 4.5, icon_color=C_RED)

    # Objectives column
    add_text(slide, "Project Objectives", 7.0, 1.5, 5.5, 0.5,
             font_size=16, bold=True, color=C_EMERALD)
    objectives = [
        "Build end-to-end ETL pipeline with SQLite star schema",
        "Compute risk-adjusted metrics (Sharpe, Sortino, Alpha, VaR)",
        "Generate 15+ EDA visualizations with business insights",
        "Create 4-page interactive Power BI dashboard",
        "Implement rule-based fund recommendation engine",
    ]
    add_bullet_box(slide, objectives, 7.0, 2.1, 5.8, 4.5, icon_color=C_EMERALD)


def slide_03_data_sources(prs):
    """Slide 3: Data Sources & Datasets."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_title_block(slide, "Data Sources & Datasets", "10 structured CSV datasets | AMFI India + mfapi.in")

    datasets = [
        ("01 Fund Master", "40 schemes, AMC, category, plan type"),
        ("02 NAV History", "51,400+ daily NAV records (2020-2025)"),
        ("03 AUM by Fund House", "140 AMCs × monthly AUM data"),
        ("04 Monthly SIP Inflows", "60 months of industry SIP data"),
        ("05 Category Inflows", "Inflow by fund category per month"),
        ("06 Folio Count", "Industry folio growth (24 months)"),
        ("07 Scheme Performance", "Sharpe, Sortino, Alpha, Beta, VaR"),
        ("08 Investor Transactions", "32,780 SIP/redemption records"),
        ("09 Portfolio Holdings", "323 stock-level equity holdings"),
        ("10 Benchmark Indices", "NIFTY 50/100/SENSEX daily OHLCV"),
    ]
    for i, (name, desc) in enumerate(datasets):
        col = i // 5
        row = i % 5
        x = 0.5 + col * 6.5
        y = 1.6 + row * 1.0
        add_rect(slide, x, y, 6.0, 0.85, C_SLATE)
        add_text(slide, name, x + 0.15, y + 0.05, 2.5, 0.4,
                 font_size=11, bold=True, color=C_EMERALD)
        add_text(slide, desc, x + 0.15, y + 0.45, 5.7, 0.35,
                 font_size=10, color=C_LIGHT)

    add_text(slide, "Sources: AMFI India · mfapi.in REST API · NSE/BSE public data | Total: ~5.5 MB",
             0.5, 7.0, 12.33, 0.4, font_size=10, color=C_DIM, align=PP_ALIGN.CENTER)


def slide_04_architecture(prs):
    """Slide 4: System Architecture & ETL Pipeline."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_title_block(slide, "System Architecture & ETL Pipeline", "Extract → Transform → Load → Analyze → Visualize")

    # Pipeline phases as connected boxes
    phases = [
        ("EXTRACT", "mfapi.in REST\n+ 10 CSVs", C_BLUE, 0.3),
        ("TRANSFORM", "Pandas clean\nFFill + Dedup", C_AMBER, 2.9),
        ("LOAD", "SQLite DB\nStar Schema", C_EMERALD, 5.5),
        ("ANALYZE", "Metrics +\nEDA + Risk", C_RED, 8.1),
        ("VISUALIZE", "Power BI\n+ Reports", C_BLUE, 10.7),
    ]
    for label, desc, color, x in phases:
        add_rect(slide, x, 1.8, 2.3, 2.0, C_SLATE)
        add_rect(slide, x, 1.8, 2.3, 0.5, color)
        add_text(slide, label, x, 1.82, 2.3, 0.45,
                 font_size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, desc, x + 0.1, 2.4, 2.1, 1.2,
                 font_size=10, color=C_LIGHT, align=PP_ALIGN.CENTER)

    # Star schema tables
    add_text(slide, "Star Schema Database (8 Tables)", 0.5, 4.2, 12, 0.5,
             font_size=14, bold=True, color=C_EMERALD)

    tables = [
        ("dim_fund", C_BLUE), ("dim_date", C_BLUE),
        ("fact_nav", C_EMERALD), ("fact_aum", C_EMERALD),
        ("fact_sip", C_EMERALD), ("fact_transactions", C_AMBER),
        ("fact_performance", C_AMBER), ("fact_portfolio", C_AMBER),
    ]
    for i, (tbl, color) in enumerate(tables):
        col = i % 4
        row = i // 4
        x = 0.5 + col * 3.2
        y = 4.9 + row * 1.0
        add_rect(slide, x, y, 3.0, 0.8, C_SLATE)
        add_rect(slide, x, y, 0.15, 0.8, color)
        add_text(slide, tbl, x + 0.25, y + 0.15, 2.7, 0.5,
                 font_size=11, bold=True, color=C_LIGHT)


def slide_05_eda(prs):
    """Slide 5: EDA Highlights."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_title_block(slide, "EDA Highlights", "15 visualizations | 10 business insights")
    insights = [
        "SIP inflows grew 3x: ₹11K Cr (Jan'23) to ₹31K Cr (Dec'25)",
        "T30 cities = 67% volume | B30 growing faster at +18% YoY",
        "36-45 age group = highest SIP investors (32% of volume)",
        "Female investors show higher SIP continuity despite lower avg amount",
        "Large Cap NAV correlation > 0.92 | Small Cap more idiosyncratic",
        "Pharma & Banking sectors dominate equity fund holdings",
        "Redemptions spike in Jan/Apr (tax planning) and Oct (corrections)",
        "UPI is now the dominant payment mode at 42% of transactions",
        "Folio count grew 18% YoY — reflecting market participation surge",
        "December is consistently the highest SIP inflow month",
    ]
    for i, insight in enumerate(insights):
        row = i % 5
        col = i // 5
        x = 0.5 + col * 6.5
        y = 1.6 + row * 1.0
        add_rect(slide, x, y, 6.0, 0.85, C_SLATE)
        add_text(slide, f"{i+1:02d}.", x + 0.1, y + 0.2, 0.6, 0.5,
                 font_size=14, bold=True, color=C_EMERALD)
        add_text(slide, insight, x + 0.7, y + 0.18, 5.2, 0.55,
                 font_size=10, color=C_LIGHT)


def slide_06_performance(prs):
    """Slide 6: Fund Performance Scorecard."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_title_block(slide, "Fund Performance Analytics", "Composite scorecard | Risk-adjusted returns")
    add_kpi_box(slide, "24.6%", "Top 3yr CAGR\n(Small Cap)", 0.5, 1.6, val_color=C_EMERALD)
    add_kpi_box(slide, "0.94", "Best Sharpe\n(SBI Small Cap)", 2.5, 1.6, val_color=C_BLUE)
    add_kpi_box(slide, "1.0%", "Avg Expense\nDirect Plans", 4.5, 1.6, val_color=C_AMBER)
    add_kpi_box(slide, "94.8", "Top Composite\nScore", 6.5, 1.6, val_color=C_RED)
    add_kpi_box(slide, "5", "Metrics in\nScorecard", 8.5, 1.6, val_color=C_EMERALD)
    add_kpi_box(slide, "40", "Funds Ranked",  10.5, 1.6, val_color=C_BLUE)
    add_text(slide, "Composite Score = 30% CAGR + 25% Sharpe + 20% Alpha + 15% Expense(inv) + 10% MaxDD(inv)",
             0.5, 3.2, 12.33, 0.5, font_size=12, color=C_DIM, italic=True, align=PP_ALIGN.CENTER)
    add_bullet_box(slide, [
        "SBI Small Cap Direct is the #1 ranked fund (Composite Score: 94.8)",
        "All top-10 funds by Sharpe Ratio are either Small Cap or Flexi Cap",
        "Direct plans outperform Regular plans by 0.6-0.9% annually (expense drag)",
        "Max Drawdown range: -10% (Low Risk) to -25% (Very High Risk)",
        "Alpha > 1% for 18 out of 40 funds — indicating active management value",
    ], 0.5, 3.8, 12, 3.2)


def slide_07_risk(prs):
    """Slide 7: Advanced Risk Metrics."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_title_block(slide, "Advanced Risk Metrics", "VaR, CVaR, Rolling Sharpe, HHI Concentration")
    # Rolling Sharpe chart
    chart_path = REPORTS_DIR / "rolling_sharpe_chart.png"
    if chart_path.exists():
        slide.shapes.add_picture(str(chart_path), Inches(0.5), Inches(1.5),
                                 Inches(6.5), Inches(4.5))
    add_text(slide, "Key Risk Findings", 7.5, 1.5, 5.3, 0.5,
             font_size=14, bold=True, color=C_EMERALD)
    add_bullet_box(slide, [
        "VaR 95%: Small Cap = -1.8% to -2.2% daily",
        "VaR 95%: Large Cap = -0.9% to -1.1% daily",
        "CVaR shows 40% worse loss on the worst 5% of days",
        "Rolling Sharpe peaked at 2.5+ in mid-2023 bull run",
        "38.2% of SIP investors are 'at-risk' (gap > 35 days)",
        "HHI: Concentration does NOT predict higher returns",
        "Most diversified fund: Nippon India Large Cap (HHI=0.08)",
    ], 7.5, 2.1, 5.5, 4.5)


def slide_08_dashboard_market(prs):
    """Slide 8: Power BI Dashboard — Market & Performance (Pages 1 & 2)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_title_block(slide, "Power BI Dashboard — Market & Performance",
                      "Pages 1-2: Industry overview, AUM trends, fund scorecard")

    # Embed page 1 and page 2 screenshots side by side
    pg1 = DASHBOARD_DIR / "page_1.png"
    pg2 = DASHBOARD_DIR / "page_2.png"
    if pg1.exists():
        slide.shapes.add_picture(str(pg1), Inches(0.3), Inches(1.5),
                                 Inches(6.3), Inches(4.2))
    if pg2.exists():
        slide.shapes.add_picture(str(pg2), Inches(6.8), Inches(1.5),
                                 Inches(6.3), Inches(4.2))

    add_text(slide, "Page 1: AUM trend, Top AMCs, SIP inflow, Folio growth",
             0.3, 5.9, 6.3, 0.5, font_size=10, color=C_DIM, align=PP_ALIGN.CENTER)
    add_text(slide, "Page 2: Risk-return scatter, NAV trends, Performance table",
             6.8, 5.9, 6.3, 0.5, font_size=10, color=C_DIM, align=PP_ALIGN.CENTER)

    add_text(slide, "Technology: Power BI Desktop | DAX Measures | Star Schema | Slicers + Drill-through",
             0.5, 6.8, 12.33, 0.4, font_size=10, color=C_DIM, align=PP_ALIGN.CENTER)


def slide_09_dashboard_investor(prs):
    """Slide 9: Power BI Dashboard — Investor & Benchmark (Pages 3 & 4)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_title_block(slide, "Power BI Dashboard — Investor & Benchmark",
                      "Pages 3-4: Demographics, geographic heatmap, benchmark comparison")

    pg3 = DASHBOARD_DIR / "page_3.png"
    pg4 = DASHBOARD_DIR / "page_4.png"
    if pg3.exists():
        slide.shapes.add_picture(str(pg3), Inches(0.3), Inches(1.5),
                                 Inches(6.3), Inches(4.2))
    if pg4.exists():
        slide.shapes.add_picture(str(pg4), Inches(6.8), Inches(1.5),
                                 Inches(6.3), Inches(4.2))

    add_text(slide, "Page 3: Age/Gender splits, Geo heatmap, City tier analysis",
             0.3, 5.9, 6.3, 0.5, font_size=10, color=C_DIM, align=PP_ALIGN.CENTER)
    add_text(slide, "Page 4: NAV vs Index, Category flows, Sector HHI",
             6.8, 5.9, 6.3, 0.5, font_size=10, color=C_DIM, align=PP_ALIGN.CENTER)

    add_text(slide, "Consistent fintech blue/purple theme | Slicers on every page | Interactive tooltips",
             0.5, 6.8, 12.33, 0.4, font_size=10, color=C_DIM, align=PP_ALIGN.CENTER)


def slide_10_recommender(prs):
    """Slide 10: Recommender Logic & Cohort Analysis."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_title_block(slide, "Fund Recommender & Cohort Analysis",
                      "Rule-based engine + investor behavior insights")

    # Recommender section (left)
    add_text(slide, "Recommendation Engine", 0.5, 1.5, 6, 0.5,
             font_size=14, bold=True, color=C_EMERALD)
    add_bullet_box(slide, [
        "Maps risk appetite → compatible fund risk grades",
        "Ranks candidates by Sharpe Ratio (top-3)",
        "Low: Low risk → Bond/Debt funds",
        "Moderate: Moderate + Moderately High → Flexi/Large Cap",
        "High: High + Very High → Small Cap/Mid Cap",
        "CLI: python scripts/recommender.py --risk Moderate",
    ], 0.5, 2.1, 5.8, 4.5, icon_color=C_BLUE)

    # Cohort section (right)
    add_text(slide, "Investor Cohort Insights", 7.0, 1.5, 6, 0.5,
             font_size=14, bold=True, color=C_EMERALD)
    add_bullet_box(slide, [
        "2022 cohort: Highest avg SIP at ₹8,200+/month",
        "38.2% of 6+ SIP investors are 'at-risk' (gap >35 days)",
        "Under-30 investors: Fastest growing segment (+28% YoY)",
        "Tier-2 cities: Lower churn (32%) vs Tier-1 (40%)",
        "Jan/Apr: Peak discontinuation months (tax season effect)",
        "HHI analysis: Concentration ≠ better returns",
    ], 7.0, 2.1, 5.8, 4.5, icon_color=C_AMBER)


def slide_11_recommendations(prs):
    """Slide 11: Strategic Recommendations."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_title_block(slide, "Strategic Recommendations", "Data-driven insights for AMCs, distributors & investors")
    recs = [
        (C_EMERALD, "For Investors",
         ["SIP investors in 2022 cohort benefited most — stay invested through cycles",
          "Small Cap: High returns (24% CAGR) but requires 3+ year horizon",
          "Prefer Direct Plans: Save 0.7-1.0% annually over Regular Plans"]),
        (C_BLUE, "For Fund Managers",
         ["SIP retention: 38% at-risk investors need proactive outreach",
          "Expense ratio reduction in Regular plans improves competitiveness",
          "Diversified portfolios (low HHI) delivered comparable risk-adjusted returns"]),
        (C_AMBER, "For Bluestock Fintech",
         ["Productionize recommender as a REST API endpoint",
          "Rolling Sharpe monitoring for fund-of-funds rebalancing",
          "Cohort analysis powers personalized investor communication"]),
    ]
    for i, (color, title, bullets) in enumerate(recs):
        x = 0.5 + i * 4.3
        add_rect(slide, x, 1.5, 4.0, 5.5, C_SLATE)
        add_rect(slide, x, 1.5, 4.0, 0.55, color)
        add_text(slide, title, x + 0.15, 1.55, 3.7, 0.45,
                 font_size=13, bold=True, color=C_WHITE)
        add_bullet_box(slide, bullets, x + 0.1, 2.2, 3.8, 4.5, icon_color=color)


def slide_12_thankyou(prs):
    """Slide 12: Conclusion & Thank You."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_rect(slide, 0, 0, 13.33, 7.5, C_DARK)
    add_rect(slide, 0, 0, 13.33, 2.0, C_SLATE)
    add_text(slide, "Thank You",
             0.5, 0.2, 12.33, 1.0, font_size=40, bold=True, color=C_EMERALD, align=PP_ALIGN.CENTER)
    add_text(slide, "Bluestock Fintech | Mutual Fund Analytics Capstone",
             0.5, 1.2, 12.33, 0.6, font_size=14, color=C_DIM, align=PP_ALIGN.CENTER)
    add_line(slide, 1.5, 2.1, 10.33)

    achievements = [
        "10 raw datasets → cleaned, modelled, and loaded into SQLite star schema",
        "15 EDA visualizations with 10 actionable business insights",
        "40 funds scored on 5 performance dimensions (CAGR, Sharpe, Alpha, Expense, MaxDD)",
        "4-page interactive Power BI dashboard with DAX-driven KPI cards",
        "Advanced risk analytics: VaR, CVaR, Rolling Sharpe, HHI, Cohort, SIP Continuity",
        "Rule-based fund recommendation engine with CLI interface",
        "Comprehensive 15–20 page final report + this 12-slide presentation",
    ]
    for i, a in enumerate(achievements):
        y = 2.3 + i * 0.65
        add_rect(slide, 0.5, y, 12.33, 0.55, C_SLATE)
        add_text(slide, "✓", 0.65, y + 0.07, 0.4, 0.4,
                 font_size=14, bold=True, color=C_EMERALD)
        add_text(slide, a, 1.1, y + 0.07, 11.5, 0.4,
                 font_size=12, color=C_LIGHT)

    add_text(slide, "Kunal Gupta  |  Internship @ Bluestock Fintech  |  July 2026",
             0.5, 7.0, 12.33, 0.45, font_size=13, color=C_DIM,
             align=PP_ALIGN.CENTER, italic=True)


# ── Main ──────────────────────────────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_cover(prs)            # 1. Cover
    slide_02_problem(prs)          # 2. Problem Statement & Objectives
    slide_03_data_sources(prs)     # 3. Data Sources & Datasets
    slide_04_architecture(prs)     # 4. System Architecture & ETL Pipeline
    slide_05_eda(prs)              # 5. EDA Highlights
    slide_06_performance(prs)      # 6. Fund Performance Scorecard
    slide_07_risk(prs)             # 7. Advanced Risk Metrics
    slide_08_dashboard_market(prs) # 8. Dashboard — Market & Performance
    slide_09_dashboard_investor(prs) # 9. Dashboard — Investor & Benchmark
    slide_10_recommender(prs)      # 10. Recommender & Cohort
    slide_11_recommendations(prs)  # 11. Strategic Recommendations
    slide_12_thankyou(prs)         # 12. Thank You

    out = REPORTS_DIR / "Presentation.pptx"
    prs.save(str(out))
    print(f"[OK] Presentation.pptx saved -> {out}")
    print(f"     Total slides: {len(prs.slides)}")
    return out


if __name__ == "__main__":
    build()
